# [ppt 수정사항]

# 부적합 추가

# 고객불만번호를 고객불만보고서로 정정: 
# 고객불만번호가 포함된 엑셀 파일은 이제 complaints_report로 지정

# 부적합번호에 대한 처리 추가: 
# 부적합번호가 포함된 엑셀 파일은 oos_report로 지정되며, 
# 추출되는 컬럼은 '부적합번호', '제목', 'QA 검토자', '발생부서', '제품', '근본 원인', '부적합품 후속 처리 계획', '조사 세부사항', '사유', '완료 요약', '후속 조치 완료 여부'

# 파일명 수정: 
# 생성되는 PPT 파일의 이름을 deviation_report.pptx, complaints_report.pptx, oos_report.pptx로 각각 변경

import dash
from dash import html, dcc
import dash_bootstrap_components as dbc
from dash.dependencies import Input, Output, State
import dataiku
import io
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import base64
import logging

logger = logging.getLogger(__name__)

#app = dash.Dash(__name__)

app.layout = html.Div([
    html.H2("엑셀 파일을 PPT 보고서로 변환하기"),
   
    # 글자 크기, 텍스트 정렬, 보고서 모드 선택을 가로로 정렬하고 Bold로 표시
    dbc.Row([
        dbc.Col([
            html.Label(html.B("글자 크기 선택:")),
            dcc.RadioItems(
                id='font_size',
                options=[
                    {'label': '크기 8', 'value': 8},
                    {'label': '크기 9', 'value': 9},
                    {'label': '크기 10', 'value': 10}
                ],
                value=8
            ),
        ], width=4),
        dbc.Col([
            html.Label(html.B("텍스트 정렬 선택:")),
            dcc.RadioItems(
                id='alignment',
                options=[
                    {'label': '좌측 정렬', 'value': 'left'},
                    {'label': '가운데 정렬', 'value': 'center'},
                    {'label': '우측 정렬', 'value': 'right'}
                ],
                value='left'
            ),
        ], width=4),
        dbc.Col([
            html.Label(html.B("보고서 모드 선택:")),
            dcc.RadioItems(
                id='mode',
                options=[
                    {'label': '라이트 모드', 'value': 'light'},
                    {'label': '다크 모드', 'value': 'dark'},
                    {'label': '여름휴가 모드', 'value': 'vacation'}
                ],
                value='light'
            ),
        ], width=4),
    ], className="mb-4"),
   
    # 엑셀 파일 업로드
    html.Label("엑셀 파일 업로드:"),
    dcc.Upload(
        id='upload_data',
        children=html.Div([
            '파일을 드래그하거나 ',
            html.A('클릭하여 선택하세요')
        ]),
        style={
            'width': '100%',
            'height': '60px',
            'lineHeight': '60px',
            'borderWidth': '1px',
            'borderStyle': 'dashed',
            'borderRadius': '5px',
            'textAlign': 'center',
            'margin': '10px'
        },
        multiple=True
    ),

    html.Div(id='output-data-upload'),

    # 버튼을 가로로 배치
    dbc.Row([
        dbc.Col(dbc.Button("PPT 생성", id="generate_ppt", color="primary"), width=2),
    ], justify="center", className="mt-3"),

    dcc.Download(id="download_ppt_iltal"),
    dcc.Download(id="download_ppt_gongmun"),
    dcc.Download(id="download_ppt_oos")
])

@app.callback(
    Output('output-data-upload', 'children'),
    Input('upload_data', 'contents'),
    State('upload_data', 'filename'),
    prevent_initial_call=True
)
def update_output(list_of_contents, list_of_names):
    if list_of_contents is None or list_of_names is None:
        raise dash.exceptions.PreventUpdate
    return html.Div([f'업로드된 파일: {", ".join(list_of_names)}'])

def create_ppt(df, selected_mode, selected_font_size, selected_alignment, pptx_title):
    prs = Presentation()

    for index, row in df.iterrows():
        slide_layout = prs.slide_layouts[5]  # 빈 슬라이드 레이아웃
        slide = prs.slides.add_slide(slide_layout)

        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = selected_mode['background']

        rows, cols = len(df.columns), 2
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(0.5 * rows)

        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        table.columns[0].width = Inches(2.0)
        table.columns[1].width = Inches(7.0)

        for i, col in enumerate(df.columns):
            cell_0 = table.cell(i, 0)
            cell_0.text = col
            cell_1 = table.cell(i, 1)
            cell_1.text = str(row[col])

            for cell in (cell_0, cell_1):
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(selected_font_size)
                        run.font.name = 'Arial'
                        run.font.color.rgb = selected_mode['text']
                    paragraph.alignment = selected_alignment
                cell.text_frame.word_wrap = True

    pptx_stream = io.BytesIO()
    prs.save(pptx_stream)
    pptx_stream.seek(0)

    return pptx_stream.getvalue()

@app.callback(
    [Output('download_ppt_iltal', 'data'), Output('download_ppt_gongmun', 'data'), Output('download_ppt_oos', 'data')],
    Input('generate_ppt', 'n_clicks'),
    State('upload_data', 'contents'),
    State('upload_data', 'filename'),
    State('font_size', 'value'),
    State('alignment', 'value'),
    State('mode', 'value'),
    prevent_initial_call=True
)
def generate_ppt(n_clicks, list_of_contents, list_of_names, font_size, alignment, mode):
    if list_of_contents is None or list_of_names is None:
        raise dash.exceptions.PreventUpdate

    iltal_dfs = []
    gongmun_dfs = []
    oos_dfs = []

    modes = {
        'light': {'background': RGBColor(255, 255, 255), 'text': RGBColor(0, 0, 0)},
        'dark': {'background': RGBColor(0, 0, 0), 'text': RGBColor(255, 255, 255)},
        'vacation': {'background': RGBColor(173, 216, 230), 'text': RGBColor(0, 0, 0)}
    }
    selected_mode = modes[mode]

    alignments = {
        'left': PP_ALIGN.LEFT,
        'center': PP_ALIGN.CENTER,
        'right': PP_ALIGN.RIGHT
    }
    selected_alignment = alignments[alignment]

    for content, name in zip(list_of_contents, list_of_names):
        content_type, content_string = content.split(',')
        decoded = base64.b64decode(content_string)
        df = pd.read_excel(io.BytesIO(decoded))

        if '일탈번호' in df.columns:
            columns_to_keep = ['일탈번호', '일탈등급', '제목', 'QA 검토자', '작성일', '일탈기준', '일탈내용', '작업자오류내용']
            iltal_dfs.append(df[columns_to_keep])
        elif '고객불만번호' in df.columns:
            columns_to_keep = ['고객불만번호', '제목', '불만발생일', '조사담당자의견', '고객요구사항']
            gongmun_dfs.append(df[columns_to_keep])
        elif '부적합번호' in df.columns:
            columns_to_keep = ['부적합번호', '제목', 'QA 검토자', '발생부서', '제품', '근본 원인', '부적합품 후속 처리 계획', '조사 세부사항', '사유', '완료 요약', '후속 조치 완료 여부']
            oos_dfs.append(df[columns_to_keep])

    if iltal_dfs:
        combined_iltal_df = pd.concat(iltal_dfs).sort_values(by='일탈번호')
        deviation_report = create_ppt(combined_iltal_df, selected_mode, font_size, selected_alignment, "일탈 보고서")
    else:
        deviation_report = None

    if gongmun_dfs:
        combined_gongmun_df = pd.concat(gongmun_dfs).sort_values(by='고객불만번호')
        complaints_report = create_ppt(combined_gongmun_df, selected_mode, font_size, selected_alignment, "고객불만 보고서")
    else:
        complaints_report = None

    if oos_dfs:
        combined_oos_df = pd.concat(oos_dfs).sort_values(by='부적합번호')
        oos_report = create_ppt(combined_oos_df, selected_mode, font_size, selected_alignment, "부적합 보고서")
    else:
        oos_report = None

    return (
        dcc.send_bytes(deviation_report, 'deviation_report.pptx') if deviation_report else None,
        dcc.send_bytes(complaints_report, 'complaints_report.pptx') if complaints_report else None,
        dcc.send_bytes(oos_report, 'oos_report.pptx') if oos_report else None
    )

if __name__ == '__main__':
    app.run_server(debug=True)




# [엑셀 파일]

# 각 유형별로 개별 시트 출력됨
# 3가지 폴더 추가함
# 나머지는 칼럼만 추가
