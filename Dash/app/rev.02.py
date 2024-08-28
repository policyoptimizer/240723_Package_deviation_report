# 엑셀 업로드 > ppt 변환 까지 가능
# 현재까지 Best!!

import dash
from dash import html, dcc
import dash_bootstrap_components as dbc
from dash.dependencies import Input, Output, State
import dataiku
import io
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
import base64
import logging

logger = logging.getLogger(__name__)

dbc_css = "https://cdn.jsdelivr.net/gh/AnnMarieW/dash-bootstrap-templates/dbc.min.css"
# 주석 처리해도 앱이 작동하도록 하는 코드
# app = dash.Dash(__name__, external_stylesheets=[dbc.themes.SUPERHERO, dbc_css])

# 웹앱 레이아웃 정의
app.layout = html.Div([
    html.H2("Excel to PPT Report Generator"),
   
    html.Label("Select Font Size:"),
    dcc.RadioItems(
        id='font_size',
        options=[
            {'label': 'Size 8', 'value': 8},
            {'label': 'Size 9', 'value': 9},
            {'label': 'Size 10', 'value': 10}
        ],
        value=8
    ),
   
    html.Label("Select Left Margin:"),
    dcc.RadioItems(
        id='left_margin',
        options=[
            {'label': '2.0 Inches', 'value': 2.0},
            {'label': '2.5 Inches', 'value': 2.5},
            {'label': '3.0 Inches', 'value': 3.0}
        ],
        value=2.0
    ),
   
    html.Label("Upload Excel File:"),
    dcc.Upload(
        id='upload_data',
        children=html.Div([
            'Drag and Drop or ',
            html.A('Select Files')
        ]),
        style={
            'width': '100%',
            'height': '60px',
            'lineHeight': '60px',
            'borderWidth': '1px',
            'borderStyle': 'dashed',
            'borderRadius': '5px',
            'textAlign': 'center',
            'margin': '10px'
        },
        multiple=False
    ),
   
    html.Div(id='output-data-upload'),
   
    dbc.Button("Generate PPT", id="generate_ppt", color="primary"),
    dcc.Download(id="download_ppt")
])

@app.callback(
    Output('output-data-upload', 'children'),
    Input('upload_data', 'contents'),
    State('upload_data', 'filename'),
    prevent_initial_call=True
)
def update_output(file_content, filename):
    if file_content is None:
        raise dash.exceptions.PreventUpdate
    return html.Div([f'Uploaded file: {filename}'])

@app.callback(
    Output('download_ppt', 'data'),
    Input('generate_ppt', 'n_clicks'),
    State('upload_data', 'contents'),
    State('upload_data', 'filename'),
    State('font_size', 'value'),
    State('left_margin', 'value'),
    prevent_initial_call=True
)
def generate_ppt(n_clicks, file_content, filename, font_size, left_margin):
    if file_content is None:
        raise dash.exceptions.PreventUpdate

    content_type, content_string = file_content.split(',')
    decoded = base64.b64decode(content_string)
    df = pd.read_excel(io.BytesIO(decoded))
   
    # 필요한 열만 남기기
    columns_to_keep = ['고객불만번호', '제목', '불만발생일', '조사담당자의견', '고객요구사항']
    df = df[columns_to_keep]
   
    # 프레젠테이션 객체 생성
    prs = Presentation()

    # 각 행에 대해 슬라이드를 생성
    for index, row in df.iterrows():
        slide_layout = prs.slide_layouts[5]  # 빈 슬라이드 레이아웃
        slide = prs.slides.add_slide(slide_layout)

        # 테이블 추가
        rows, cols = len(df.columns), 2
        left = Inches(left_margin)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(0.5 * rows)  # 행의 높이를 줄임

        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # 테이블 첫 열의 너비 조정
        table.columns[0].width = Inches(2.0)
        table.columns[1].width = Inches(7.0)

        # 테이블 내용 채우기
        for i, col in enumerate(df.columns):
            cell_0 = table.cell(i, 0)
            cell_0.text = col
            cell_1 = table.cell(i, 1)
            cell_1.text = str(row[col])

            # 글자 크기 및 텍스트 프레임 설정
            for cell in (cell_0, cell_1):
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(font_size)
                        run.font.name = 'Arial'
                cell.text_frame.word_wrap = True

    # PPT 파일을 바이트 스트림으로 변환
    pptx_stream = io.BytesIO()
    prs.save(pptx_stream)
    pptx_stream.seek(0)
   
    return dcc.send_bytes(pptx_stream.getvalue(), f'{filename.rsplit(".", 1)[0]}_report.pptx')

# Dataiku Dash 웹앱 실행
if __name__ == '__main__':
    app.run_server(debug=True)

