# 이걸로 메일 송부했음. 중간공유.
# 현재까지 best

import dash
from dash import html, dcc
import dash_bootstrap_components as dbc
from dash.dependencies import Input, Output, State
import dataiku
import io
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import base64
import logging

logger = logging.getLogger(__name__)

#app = dash.Dash(__name__)

app.layout = html.Div([
    html.H2("엑셀 파일을 PPT 보고서로 변환하기"),
   
    # 글자 크기, 텍스트 정렬, 보고서 모드 선택을 가로로 정렬하고 Bold로 표시
    dbc.Row([
        dbc.Col([
            html.Label(html.B("글자 크기 선택:")),
            dcc.RadioItems(
                id='font_size',
                options=[
                    {'label': '크기 8', 'value': 8},
                    {'label': '크기 9', 'value': 9},
                    {'label': '크기 10', 'value': 10}
                ],
                value=8
            ),
        ], width=4),
        dbc.Col([
            html.Label(html.B("텍스트 정렬 선택:")),
            dcc.RadioItems(
                id='alignment',
                options=[
                    {'label': '좌측 정렬', 'value': 'left'},
                    {'label': '가운데 정렬', 'value': 'center'},
                    {'label': '우측 정렬', 'value': 'right'}
                ],
                value='left'
            ),
        ], width=4),
        dbc.Col([
            html.Label(html.B("보고서 모드 선택:")),
            dcc.RadioItems(
                id='mode',
                options=[
                    {'label': '라이트 모드', 'value': 'light'},
                    {'label': '다크 모드', 'value': 'dark'},
                    {'label': '여름휴가 모드', 'value': 'vacation'}
                ],
                value='light'
            ),
        ], width=4),
    ], className="mb-4"),
   
    # 엑셀 파일 업로드
    html.Label("엑셀 파일 업로드:"),
    dcc.Upload(
        id='upload_data',
        children=html.Div([
            '파일을 드래그하거나 ',
            html.A('클릭하여 선택하세요')
        ]),
        style={
            'width': '100%',
            'height': '60px',
            'lineHeight': '60px',
            'borderWidth': '1px',
            'borderStyle': 'dashed',
            'borderRadius': '5px',
            'textAlign': 'center',
            'margin': '10px'
        },
        multiple=True
    ),

    html.Div(id='output-data-upload'),

    # 버튼을 가로로 배치
    dbc.Row([
        dbc.Col(dbc.Button("PPT 생성", id="generate_ppt", color="primary"), width=2),
    ], justify="center", className="mt-3"),

    dcc.Download(id="download_ppt_iltal"),
    dcc.Download(id="download_ppt_gongmun")
])

@app.callback(
    Output('output-data-upload', 'children'),
    Input('upload_data', 'contents'),
    State('upload_data', 'filename'),
    prevent_initial_call=True
)
def update_output(list_of_contents, list_of_names):
    if list_of_contents is None or list_of_names is None:
        raise dash.exceptions.PreventUpdate
    return html.Div([f'업로드된 파일: {", ".join(list_of_names)}'])

def create_ppt(df, selected_mode, selected_font_size, selected_alignment, pptx_title):
    prs = Presentation()

    for index, row in df.iterrows():
        slide_layout = prs.slide_layouts[5]  # 빈 슬라이드 레이아웃
        slide = prs.slides.add_slide(slide_layout)

        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = selected_mode['background']

        rows, cols = len(df.columns), 2
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(0.5 * rows)

        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        table.columns[0].width = Inches(2.0)
        table.columns[1].width = Inches(7.0)

        for i, col in enumerate(df.columns):
            cell_0 = table.cell(i, 0)
            cell_0.text = col
            cell_1 = table.cell(i, 1)
            cell_1.text = str(row[col])

            for cell in (cell_0, cell_1):
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(selected_font_size)
                        run.font.name = 'Arial'
                        run.font.color.rgb = selected_mode['text']
                    paragraph.alignment = selected_alignment
                cell.text_frame.word_wrap = True

    pptx_stream = io.BytesIO()
    prs.save(pptx_stream)
    pptx_stream.seek(0)

    return pptx_stream.getvalue()

@app.callback(
    [Output('download_ppt_iltal', 'data'), Output('download_ppt_gongmun', 'data')],
    Input('generate_ppt', 'n_clicks'),
    State('upload_data', 'contents'),
    State('upload_data', 'filename'),
    State('font_size', 'value'),
    State('alignment', 'value'),
    State('mode', 'value'),
    prevent_initial_call=True
)
def generate_ppt(n_clicks, list_of_contents, list_of_names, font_size, alignment, mode):
    if list_of_contents is None or list_of_names is None:
        raise dash.exceptions.PreventUpdate

    iltal_dfs = []
    gongmun_dfs = []

    modes = {
        'light': {'background': RGBColor(255, 255, 255), 'text': RGBColor(0, 0, 0)},
        'dark': {'background': RGBColor(0, 0, 0), 'text': RGBColor(255, 255, 255)},
        'vacation': {'background': RGBColor(173, 216, 230), 'text': RGBColor(0, 0, 0)}
    }
    selected_mode = modes[mode]

    alignments = {
        'left': PP_ALIGN.LEFT,
        'center': PP_ALIGN.CENTER,
        'right': PP_ALIGN.RIGHT
    }
    selected_alignment = alignments[alignment]

    for content, name in zip(list_of_contents, list_of_names):
        content_type, content_string = content.split(',')
        decoded = base64.b64decode(content_string)
        df = pd.read_excel(io.BytesIO(decoded))

        if '일탈번호' in df.columns:
            columns_to_keep = ['일탈번호', '일탈등급', '제목', 'QA 검토자', '작성일', '일탈기준', '일탈내용', '작업자오류내용']
            iltal_dfs.append(df[columns_to_keep])
        elif '고객불만번호' in df.columns:
            columns_to_keep = ['고객불만번호', '제목', '불만발생일', '조사담당자의견', '고객요구사항']
            gongmun_dfs.append(df[columns_to_keep])

    if iltal_dfs:
        combined_iltal_df = pd.concat(iltal_dfs).sort_values(by='일탈번호')
        iltal_ppt = create_ppt(combined_iltal_df, selected_mode, font_size, selected_alignment, "일탈 보고서")
    else:
        iltal_ppt = None

    if gongmun_dfs:
        combined_gongmun_df = pd.concat(gongmun_dfs).sort_values(by='고객불만번호')
        gongmun_ppt = create_ppt(combined_gongmun_df, selected_mode, font_size, selected_alignment, "부적합 보고서")
    else:
        gongmun_ppt = None

    return (
        dcc.send_bytes(iltal_ppt, 'iltal_report.pptx') if iltal_ppt else None,
        dcc.send_bytes(gongmun_ppt, 'gongmun_report.pptx') if gongmun_ppt else None
    )

if __name__ == '__main__':
    app.run_server(debug=True)
