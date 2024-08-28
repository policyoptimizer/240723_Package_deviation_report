# 큰 틀에서 파일 업로드 및 처리 가능

# Import necessary libraries
import dash
from dash import html
from dash import dcc
import dash_bootstrap_components as dbc
from dash.dependencies import Input, Output, State, ALL, MATCH
from dash.exceptions import PreventUpdate

from webapps.utils import get_managed_folder_list, get_files_in_folder

import logging
import dataiku
import io
import base64
import pandas as pd
from pptx import Presentation
from pptx.util import Inches

logger = logging.getLogger(__name__)

dbc_css = "https://cdn.jsdelivr.net/gh/AnnMarieW/dash-bootstrap-templates/dbc.min.css"
# use the style of examples on the Plotly documentation
# app = dash.Dash(__name__, external_stylesheets=[dbc.themes.SUPERHERO, dbc_css])  # 주석 처리된 부분

home_layout = html.Div([
    dcc.Store(id='sequential_call'),
    dbc.Row([html.H2("Select a managed folder"), ]),
    dbc.Row([dcc.Dropdown(id='select_managed', options=[], placeholder="Select a managed folder", className="dbc"), ]),
    dbc.Row([html.H3("Upload a file to the managed folder"), ]),
    dbc.Row(
        dcc.Upload(
            id='upload_data',
            children=html.Div(id='output-data-upload', children=
            ["Drag and drop or click to select a file to upload."]
                              ),
            style={
                "height": "60px",
                "lineHeight": "60px",
                "borderWidth": "1px",
                "borderStyle": "dashed",
                "borderRadius": "5px",
                "textAlign": "center",
                "margin": "10px",
            },
            multiple=True,
            disabled=True
        ),
    ),
    dbc.Row([html.H3("Files from the managed folder"), ]),
    dbc.Row([html.Ul(id='file_list')], className="dbc container-fluid mt-3"),
], className="container-fluid mt-3")

# build your Dash app
app.layout = home_layout

def make_download_button(filename, index):
    download_area = dcc.Download(id={'index': index, 'type': 'dld'}, data={'base64': True})
    button = html.Button(filename, id={'index': index, 'type': 'btn', 'filename': filename})
    layout = html.Li(html.Div(children=[button, download_area]))
    return layout

@app.callback(
    Output('select_managed', 'options'),
    Input('select_managed', 'id')
)
def load_select(_):
    ids_and_names = get_managed_folder_list()
    return [{'value': ian[0], 'label': ian[1]} for ian in ids_and_names]

@app.callback(
    [
        Output('upload_data', 'disabled', allow_duplicate=True),
        Output('sequential_call', 'data', allow_duplicate=True),
        Output('file_list', 'children', allow_duplicate=True),
    ],
    Input('select_managed', 'value'),
    State('sequential_call', 'data'),
    prevent_initial_call=True
)
def clear_list(_, data):
    value = data or {'update_list': 0}
    value['update_list'] = value['update_list'] + 1
    return [False, value, []]

@app.callback(
    [
        Output('file_list', 'children', allow_duplicate=True),
        Output('upload_data', 'disabled', allow_duplicate=True),
    ],
    Input('sequential_call', 'data'),
    State('select_managed', 'value'),
    prevent_initial_call=True
)
def update_list(_, folder_id):
    if folder_id:
        files = get_files_in_folder(folder_id)
        if len(files) == 0:
            return [[html.Li("No file in the selected folder")], False]
        else:
            return [[make_download_button(filename, x) for x, filename in enumerate(files)], False]
    else:
        return dash.no_update, True

@app.callback(
    Output('select_managed', 'value', allow_duplicate=True),
    [Input('upload_data', 'filename'),
     Input('upload_data', 'contents')],
    State('select_managed', 'value'),
    prevent_initial_call=True
)
def upload_and_generate_ppt(uploaded_filenames, uploaded_file_contents, folder_id):
    if folder_id is not None:
        mf = dataiku.Folder(folder_id)
        if uploaded_filenames is not None and uploaded_file_contents is not None:
            for name, data in zip(uploaded_filenames, uploaded_file_contents):
                try:
                    content_type, content_string = data.split(',')
                    stream_d = base64.b64decode(content_string)
                    stream = io.BytesIO(stream_d)

                    # 엑셀 파일을 DataFrame으로 읽기
                    df = pd.read_excel(stream)
                    logger.info(f"엑셀 파일 '{name}'을 성공적으로 읽었습니다.")
                   
                    # PPT로 변환
                    ppt_stream = generate_ppt_from_excel(df)
                    ppt_filename = name.rsplit('.', 1)[0] + '.pptx'
                    logger.info(f"PPT 파일 '{ppt_filename}'을 생성 중입니다.")
                   
                    # Managed Folder에 업로드
                    mf.upload_stream(ppt_filename, ppt_stream)
                    logger.info(f"PPT 파일 '{ppt_filename}'을 Managed Folder에 업로드했습니다.")

                except Exception as e:
                    logger.error(f"파일 처리 중 오류 발생: {e}")
                    return dash.no_update

        return folder_id
    else:
        return dash.no_update

@app.callback(
    Output({'type': 'dld', 'index': MATCH}, 'data'),
    Input({'type': 'btn', 'index': MATCH, 'filename': ALL}, 'n_clicks'),
    State({'type': 'btn', 'index': MATCH, 'filename': ALL}, 'id'),
    State('select_managed', 'value'),
    prevent_initial_call=True
)
def download_file(_, id, managed):
    if id:
        mf = dataiku.Folder(managed)
        ppt_filename = id[0].get('filename', '')
        logger.info(f"다운로드 시도: {ppt_filename}")
       
        try:
            def write_file(bytes_io):
                stream = mf.get_download_stream(ppt_filename)
                bytes_io.write(stream.read())
            return dcc.send_bytes(write_file, ppt_filename)
        except Exception as e:
            logger.error(f"파일 다운로드 중 오류 발생: {e}")
            return dash.no_update
    else:
        return dash.no_update

def generate_ppt_from_excel(df):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # 데이터에서 PPT로 테이블 추가
    try:
        table = slide.shapes.add_table(len(df) + 1, len(df.columns), Inches(0.5), Inches(0.5), Inches(9), Inches(5)).table
        for col_idx, col_name in enumerate(df.columns):
            table.cell(0, col_idx).text = col_name

        for row_idx, row in df.iterrows():
            for col_idx, value in enumerate(row):
                table.cell(row_idx + 1, col_idx).text = str(value)
    except Exception as e:
        logger.error(f"PPT 테이블 생성 중 오류 발생: {e}")
        raise

    pptx_blob = io.BytesIO()
    prs.save(pptx_blob)
    pptx_blob.seek(0)
    return pptx_blob

