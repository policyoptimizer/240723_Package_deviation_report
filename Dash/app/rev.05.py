# 240829 이걸로 중간 발표 했음
# 현재까지 Best

import dash
from dash import html, dcc
import dash_bootstrap_components as dbc
from dash.dependencies import Input, Output, State
import dataiku
import io
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import base64
import logging

logger = logging.getLogger(__name__)

# app = dash.Dash(__name__, external_stylesheets=[dbc.themes.SUPERHERO, dbc_css])

app.layout = html.Div([
    html.H2("엑셀 파일을 PPT 보고서로 변환하기"),
   
    # 글자 크기, 텍스트 정렬, 보고서 모드 선택을 가로로 정렬하고 Bold로 표시
    dbc.Row([
        dbc.Col([
            html.Label(html.B("글자 크기 선택:")),
            dcc.RadioItems(
                id='font_size',
                options=[
                    {'label': '크기 8', 'value': 8},
                    {'label': '크기 9', 'value': 9},
                    {'label': '크기 10', 'value': 10}
                ],
                value=8
            ),
        ], width=4),
    ], className="mb-4"),
   
    dbc.Row([
        dbc.Col([
            html.Label(html.B("텍스트 정렬 선택:")),
            dcc.RadioItems(
                id='alignment',
                options=[
                    {'label': '좌측 정렬', 'value': 'left'},
                    {'label': '가운데 정렬', 'value': 'center'},
                    {'label': '우측 정렬', 'value': 'right'}
                ],
                value='left'
            ),
        ], width=4),
    ], className="mb-4"),
   
    dbc.Row([
        dbc.Col([
            html.Label(html.B("보고서 모드 선택:")),
            dcc.RadioItems(
                id='mode',
                options=[
                    {'label': '라이트 모드', 'value': 'light'},
                    {'label': '다크 모드', 'value': 'dark'},
                    {'label': '여름휴가 모드', 'value': 'vacation'}
                ],
                value='light'
            ),
        ], width=4),
    ], className="mb-4"),
   
    # 엑셀 파일 업로드
    html.Label("엑셀 파일 업로드:"),
    dcc.Upload(
        id='upload_data',
        children=html.Div([
            '파일을 드래그하거나 ',
            html.A('클릭하여 선택하세요')
        ]),
        style={
            'width': '100%',
            'height': '60px',
            'lineHeight': '60px',
            'borderWidth': '1px',
            'borderStyle': 'dashed',
            'borderRadius': '5px',
            'textAlign': 'center',
            'margin': '10px'
        },
        multiple=False
    ),

    html.Div(id='output-data-upload'),

    # 버튼을 가로로 배치
    dbc.Row([
        dbc.Col(dbc.Button("PPT 생성", id="generate_ppt", color="primary"), width=2),
    ], justify="center", className="mt-3"),

    dcc.Download(id="download_ppt")
])

@app.callback(
    Output('output-data-upload', 'children'),
    Input('upload_data', 'contents'),
    State('upload_data', 'filename'),
    prevent_initial_call=True
)
def update_output(file_content, filename):
    if file_content is None:
        raise dash.exceptions.PreventUpdate
    return html.Div([f'업로드된 파일: {filename}'])

@app.callback(
    Output('download_ppt', 'data'),
    Input('generate_ppt', 'n_clicks'),
    State('upload_data', 'contents'),
    State('upload_data', 'filename'),
    State('font_size', 'value'),
    State('alignment', 'value'),
    State('mode', 'value'),
    prevent_initial_call=True
)
def generate_ppt(n_clicks, file_content, filename, font_size, alignment, mode):
    if file_content is None:
        raise dash.exceptions.PreventUpdate

    content_type, content_string = file_content.split(',')
    decoded = base64.b64decode(content_string)
    df = pd.read_excel(io.BytesIO(decoded))

    # 필요한 열만 남기기
    columns_to_keep = ['고객불만번호', '제목', '불만발생일', '조사담당자의견', '고객요구사항']
    df = df[columns_to_keep]

    # 모드에 따른 설정
    modes = {
        'light': {'background': RGBColor(255, 255, 255), 'text': RGBColor(0, 0, 0)},
        'dark': {'background': RGBColor(0, 0, 0), 'text': RGBColor(255, 255, 255)},
        'vacation': {'background': RGBColor(173, 216, 230), 'text': RGBColor(0, 0, 0)}
    }
    selected_mode = modes[mode]

    # 정렬에 따른 설정
    alignments = {
        'left': PP_ALIGN.LEFT,
        'center': PP_ALIGN.CENTER,
        'right': PP_ALIGN.RIGHT
    }
    selected_alignment = alignments[alignment]

    # 프레젠테이션 객체 생성
    prs = Presentation()

    # 각 행에 대해 슬라이드를 생성
    for index, row in df.iterrows():
        slide_layout = prs.slide_layouts[5]  # 빈 슬라이드 레이아웃
        slide = prs.slides.add_slide(slide_layout)

        # 슬라이드 배경 색상 설정
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = selected_mode['background']

        # 테이블 추가
        rows, cols = len(df.columns), 2
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(0.5 * rows)

        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # 테이블 첫 열의 너비 조정
        table.columns[0].width = Inches(2.0)
        table.columns[1].width = Inches(7.0)

        # 테이블 내용 채우기
        for i, col in enumerate(df.columns):
            cell_0 = table.cell(i, 0)
            cell_0.text = col
            cell_1 = table.cell(i, 1)
            cell_1.text = str(row[col])

            # 글자 크기, 정렬 및 텍스트 프레임 설정
            for cell in (cell_0, cell_1):
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(font_size)  # 선택한 글자 크기 적용
                        run.font.name = 'Arial'
                        run.font.color.rgb = selected_mode['text']  # 선택한 텍스트 색상 적용
                    paragraph.alignment = selected_alignment  # 선택한 정렬 기준 적용
                cell.text_frame.word_wrap = True

    # PPT 파일을 바이트 스트림으로 변환
    pptx_stream = io.BytesIO()
    prs.save(pptx_stream)
    pptx_stream.seek(0)

    return dcc.send_bytes(pptx_stream.getvalue(), f'{filename.rsplit(".", 1)[0]}_report.pptx')

# Dataiku Dash 웹앱 실행
if __name__ == '__main__':
    app.run_server(debug=True)

