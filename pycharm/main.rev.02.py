# 글자 사이즈 조정, 행 높이 줄이기

import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# 경로 설정
input_file_path = r'D:\#.Secure Work Folder\BIG\Toy\24~28Y\240723 표시자재 관련 일탈사항 보고서 자동화\dat\IQMS_rev.01.xlsx'
output_file_path = r'D:\#.Secure Work Folder\BIG\Toy\24~28Y\240723 표시자재 관련 일탈사항 보고서 자동화\out\customer_complaints_report.pptx'

# 엑셀 파일 읽기
df = pd.read_excel(input_file_path)

# 필요한 열만 남기기
columns_to_keep = ['고객불만번호', '제목', '불만발생일', '조사담당자의견', '고객요구사항']
df = df[columns_to_keep]

# 프레젠테이션 객체 생성
prs = Presentation()

# 각 행에 대해 슬라이드를 생성
for index, row in df.iterrows():
    slide_layout = prs.slide_layouts[5]  # 빈 슬라이드 레이아웃
    slide = prs.slides.add_slide(slide_layout)

    # 테이블 추가
    rows, cols = len(df.columns), 2
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(0.5 * rows)  # 행의 높이를 줄임

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # 테이블 첫 열의 너비 조정
    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(7.0)

    # 테이블 내용 채우기
    for i, col in enumerate(df.columns):
        cell_0 = table.cell(i, 0)
        cell_0.text = col
        cell_1 = table.cell(i, 1)
        cell_1.text = str(row[col])

        # 글자 크기 및 텍스트 프레임 설정
        for cell in (cell_0, cell_1):
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(8)  # 글자 크기를 줄임
                    run.font.name = 'Arial'
            cell.text_frame.word_wrap = True

        # 조사담당자의견의 텍스트 길이에 따른 행 높이 조정
        if col == '조사담당자의견':
            cell_1.margin_top = Pt(2)  # 테두리 간격을 줄임
            cell_1.margin_bottom = Pt(2)

# PPT 파일 저장
prs.save(output_file_path)

print(f"PPT 보고서가 '{output_file_path}' 경로에 성공적으로 저장되었습니다.")

