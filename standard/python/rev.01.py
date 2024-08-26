# pptx 코드 환경 없어서 못하고 있음

import dataiku
from flask import request, send_file
import os
import io
from pptx import Presentation
from pptx.util import Inches, Pt

@app.route('/upload-to-dss', methods=['POST'])
def upload_to_dss():
    # 엑셀 파일 업로드 처리
    f = request.files.get('file')
    if not f:
        return "No file uploaded", 400

    # Managed Folder에 파일 저장
    mf = dataiku.Folder('box')  # 'box'는 Managed Folder의 이름입니다.
    target_path = f.filename
    mf.upload_stream(target_path, f)

    # 엑셀 파일을 기반으로 PPT 생성
    pptx_blob = generate_ppt_from_excel(mf, target_path)

    # PPT 파일을 임시 디렉토리에 저장
    with open('/tmp/customer_complaints_report.pptx', 'wb') as out_file:
        out_file.write(pptx_blob.getvalue())

    return {"status": "ok"}, 200

@app.route('/download-ppt', methods=['GET'])
def download_ppt():
    # 생성된 PPT 파일을 다운로드
    return send_file('/tmp/customer_complaints_report.pptx', as_attachment=True)

def generate_ppt_from_excel(mf, path):
    # 엑셀 파일에서 데이터를 읽어와서 PPT 생성 로직 구현
    with mf.get_download_stream(path) as stream:
        df = pd.read_excel(stream)

    # PPT 생성 로직
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    # 예시로, 데이터 첫 행을 제목으로, 나머지 행을 내용으로 하는 테이블 추가
    table = slide.shapes.add_table(len(df) + 1, len(df.columns), Inches(0.5), Inches(0.5), Inches(9), Inches(5)).table
    for col_idx, col_name in enumerate(df.columns):
        table.cell(0, col_idx).text = col_name

    for row_idx, row in df.iterrows():
        for col_idx, value in enumerate(row):
            table.cell(row_idx + 1, col_idx).text = str(value)

    # PPT 파일을 바이트로 반환
    pptx_blob = io.BytesIO()
    prs.save(pptx_blob)
    pptx_blob.seek(0)
    return pptx_blob

